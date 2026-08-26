"""Deck theme and layout helpers, keyed to the course logo.

Why python-pptx and not pptxgenjs: this repository is a Python toolchain and
the instructor is a Python programmer. A deck you can read and modify beats a
deck someone else can generate, and macOS does not ship Node.

Every deck is a Python module that calls these helpers. The helpers enforce the
things that are actually load-bearing -- consistent signalling, dark slides as
mode boundaries, extracted terms rather than prose -- so a new deck is short and
cannot drift from the house style. See docs/lecture-design.md for why the
structure is what it is.
"""
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent

# --- palette, sampled from the course logo ---------------------------------
TEAL = RGBColor(0x0E, 0x4F, 0x57)
GREEN = RGBColor(0x1A, 0x4D, 0x33)
MINT = RGBColor(0xA5, 0xD6, 0xA7)
CYAN = RGBColor(0x4F, 0xD1, 0xC5)
SILVER = RGBColor(0xA8, 0xC8, 0xD8)
INK = RGBColor(0x0B, 0x3A, 0x3F)
BODY = RGBColor(0x23, 0x42, 0x3F)
MUTED = RGBColor(0x6E, 0x8B, 0x87)
AMBER = RGBColor(0xD9, 0x8E, 0x32)
RED = RGBColor(0xB3, 0x26, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD3, 0xDE, 0xDA)
WASH = RGBColor(0xE8, 0xF4, 0xF0)

HEAD, TEXT = "Cambria", "Calibri"
W, H, M = 13.333, 7.5, 0.7          # slide size and margin, inches

sys.path.insert(0, str(ROOT))
from tools.schedule import readings_spec, sessions  # noqa: E402

PAPER_FIGURES = ROOT / "private" / "paper-figures"
PAPER_MOVIES = ROOT / "private" / "paper-movies"
MANIFEST = ROOT / "decks" / "paper_figures.yaml"


def _resolve_paper_figure(key):
    """Path to `private/paper-figures/<key>.png`, deriving it if it is a crop.

    A multi-panel journal figure dropped whole onto a slide is unreadable from
    the back of the room. The fix is to show one panel -- but a hand-cropped PNG
    is an untracked artifact that nobody can regenerate. So crops are declared
    in the manifest:

        gardner2000_fig5a:
          derived_from: gardner2000_fig5
          crop: [40, 10, 1460, 820]        # left, top, right, bottom in px

    and cut here, from the parent, on demand. The cropped file is written next
    to the parent (also gitignored) and reused until the parent changes.
    """
    img = PAPER_FIGURES / f"{key}.png"
    if img.exists():
        return img

    if not MANIFEST.exists():
        return img
    import yaml
    spec = (yaml.safe_load(MANIFEST.read_text()) or {}).get(key) or {}
    parent_key, box = spec.get("derived_from"), spec.get("crop")
    if not (parent_key and box):
        return img
    parent = PAPER_FIGURES / f"{parent_key}.png"
    if not parent.exists():
        return img

    from PIL import Image as _Image
    with _Image.open(parent) as im:
        im.crop(tuple(box)).save(img)
    return img


def _extract_poster(video, dst):
    """A still from `video`, cached at `dst`, for the slide to show at rest.

    Without one, python-pptx gives the video a generic loudspeaker icon: a deck
    full of grey speaker glyphs, which looks broken and tells the room nothing
    if the clip never gets played. A frame from ~40% through is almost always
    more representative than frame zero, which is often black or a title card.

    Silently returns the missing path if ffmpeg is absent -- the deck still
    builds, it just gets the icon.
    """
    import shutil
    import subprocess
    if not shutil.which("ffmpeg"):
        return dst
    dst.parent.mkdir(parents=True, exist_ok=True)
    dur = subprocess.run(
        ["ffprobe", "-v", "error", "-show_entries", "format=duration",
         "-of", "default=nw=1:nk=1", str(video)],
        capture_output=True, text=True).stdout.strip()
    seek = f"{float(dur) * 0.4:.2f}" if dur else "1"
    subprocess.run(["ffmpeg", "-loglevel", "error", "-y", "-ss", seek,
                    "-i", str(video), "-frames:v", "1", str(dst)],
                   capture_output=True)
    return dst


def short_cite(r):
    """What goes on a slide: 'Basu et al., PNAS 2004'.

    Not the full citation. A projected reference the room reads for two seconds
    needs to be identifiable, not complete -- the complete one is in
    docs/readings.md and on bCourses, and putting it on the slide costs two
    lines of wrap and buys nothing.
    """
    if r.get("short"):
        return r["short"]
    cite = r.get("cite", "")
    author = cite.split(",")[0].strip() or r.get("key", "?")
    year = re.search(r"\((?:19|20)\d{2}\)", cite)
    journal = re.search(r"(Nature Reviews [A-Z][a-z]+|Nature [A-Z][a-z]+|"
                        r"Nature|Science|Cell|PNAS|eLife|Mol\. Syst\. Biol\.)", cite)
    bits = [author + (" et al." if "et al" in cite or cite.count(",") > 3 else "")]
    if journal:
        bits.append(journal.group(0))
    if year:
        bits.append(year.group(0).strip("()"))
    return " ".join(bits[:1]) + (", " + " ".join(bits[1:]) if len(bits) > 1 else "")


def first_sentence(txt, limit=105):
    """The first sentence of `focus`, for the slide. The rest is in the docs."""
    s = " ".join((txt or "").split())
    if not s:
        return ""
    head = re.split(r"(?<=[.!?])\s", s)[0]
    return head if len(head) <= limit else head[:limit - 1].rstrip() + "\u2026"


class Deck:
    """A lecture deck. Every method returns self or a slide, so decks read
    top to bottom like the lecture does."""

    def __init__(self, title, session=None, subtitle="", meta=""):
        self.prs = Presentation()
        self.prs.slide_width = Inches(W)
        self.prs.slide_height = Inches(H)
        self.title_text = title
        self.subtitle = subtitle
        self.meta = meta
        self.session = session
        self.missing_figures = []
        self.loose_slots = []
        self.assignment_rendered = False
        self.assignment_overflow = None
        self.unattributed_figures = []
        self.segments = []
        self.missing_movies = []
        # Every file this deck actually read while building. Not a guess from
        # the source: the build records what it touched, so tools/manifest.py
        # can later tell whether the .pptx on the teaching machine is older
        # than any of them. See tools/manifest.py for why this is not mtimes.
        self.assets = []

    # -- the calendar, from course.yaml rather than from a typed string ------
    @property
    def date_line(self):
        """'BioE 147 / 247 · Thursday 24 September 2026 · Dwinelle 219'.

        Typed by hand this is a date that will be wrong the first time the
        calendar moves and right-looking enough that nobody checks it.
        """
        from tools.schedule import course, session_date
        c = course()
        d = session_date(self.session)
        day = d.strftime("%A %d %B %Y").replace(" 0", " ")
        return f"BioE 147 / 247  ·  {day}  ·  {c['room']}"

    def _used(self, path):
        """Record a file the build read. Idempotent; order does not matter."""
        p = Path(path)
        if not p.is_absolute():
            p = ROOT / p
        if p.exists() and p not in self.assets:
            self.assets.append(p)
        return p

    # -- slide construction --------------------------------------------------
    def _slide(self, dark):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])   # blank
        bg = ROOT / "decks" / ("bg_dark.png" if dark else "bg_light.png")
        if bg.exists():
            self._used(bg)
            pic = s.shapes.add_picture(str(bg), 0, 0,
                                       width=Inches(W), height=Inches(H))
            s.shapes._spTree.remove(pic._element)
            s.shapes._spTree.insert(2, pic._element)
        s._posb_dark = dark
        return s

    def dark(self):
        return self._slide(True)

    def light(self):
        return self._slide(False)

    # -- text ----------------------------------------------------------------
    def text(self, s, txt, x, y, w, h, size=14, font=None, bold=False,
             italic=False, color=None, align="l", valign=None, spacing=None,
             space_pt=0):
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        if valign:
            tf.vertical_anchor = {"m": MSO_ANCHOR.MIDDLE,
                                  "t": MSO_ANCHOR.TOP}[valign]
        for i, line in enumerate(str(txt).split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                           "r": PP_ALIGN.RIGHT}[align]
            if spacing:
                p.line_spacing = spacing
            if space_pt:
                p.space_after = Pt(space_pt)
            r = p.add_run()
            r.text = line
            f = r.font
            f.name = font or TEXT
            f.size = Pt(size)
            f.bold = bold
            f.italic = italic
            f.color.rgb = color or BODY
        return box

    def shape(self, s, kind, x, y, w, h, fill=None, line=None, lw=1.5):
        shp = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = Pt(lw)
        shp.shadow.inherit = False
        if shp.has_text_frame:
            shp.text_frame.text = ""
        return shp

    def image(self, s, path, x, y, w=None, h=None):
        """Place an image, **preserving its aspect ratio**.

        Passing both width and height to python-pptx stretches the image to fill
        them. So when both are given we treat them as a bounding box: scale to
        fit inside, then centre. A distorted figure from a published paper is
        both ugly and a misrepresentation of someone else's data.

        Returns the rectangle the image actually occupies, `(x, y, w, h)`, which
        is not the box you asked for whenever the aspect ratios differ. Callers
        that put a caption underneath need the real rectangle, or the caption
        floats in the empty part of the box.
        """
        p = self._used(path)
        if w is not None and h is not None:
            from PIL import Image as _Image
            iw, ih = _Image.open(p).size
            scale = min(w / iw, h / ih)
            dw, dh = iw * scale, ih * scale
            x += (w - dw) / 2
            y += (h - dh) / 2
            w, h = dw, dh
        pic = s.shapes.add_picture(str(p), Inches(x), Inches(y),
                                   Inches(w) if w else None,
                                   Inches(h) if h else None)
        return (pic.left / 914400, pic.top / 914400,
                pic.width / 914400, pic.height / 914400)

    # -- repeating layout elements ------------------------------------------
    # Words in a segment label that mean the students are working, not
    # watching. During those minutes the slide is deliberately static, so they
    # are excluded from the pacing check below.
    ACTIVITY_WORDS = ("group", "vote", "argue", "diagnostic", "laptops",
                      "worked set", "retrieval", "notes closed", "pause",
                      "i will not say", "in writing", "handout")

    # Board work is a third category, not a kind of exposition and not a kind
    # of student activity. During a derivation at the board the slide is static
    # for the same reason it is static during a vote -- it is not the medium --
    # so counting those minutes against the slide count flags a session that is
    # doing exactly the right thing.
    #
    # This is an obvious loophole, so it is a narrow and VISIBLE one: the
    # segment label has to say "board", which is a declaration in the deck
    # source, and pacing() reports board minutes as their own line rather than
    # folding them into student time. A deck that claims 40 minutes at the
    # board is making a claim you can read and disbelieve.
    BOARD_WORDS = ("board",)

    def header(self, s, badge, label):
        """Time badge + segment label. Same look every time: this is signalling,
        and signalling is one of the few interventions that helps novices and
        experts equally.

        The badge is also the only machine-readable record of how long each
        segment lasts, so it is what `pacing()` measures.
        """
        self.segments.append((badge, label))
        dark = s._posb_dark
        self.shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, 0.42, 1.62, 0.34,
                   fill=CYAN if dark else TEAL, line=None)
        self.text(s, badge, M, 0.47, 1.62, 0.3, size=11, bold=True,
                  color=INK if dark else WHITE, align="c")
        self.text(s, label.upper(), M + 1.8, 0.47, 9.5, 0.3, size=11, bold=True,
                  color=MINT if dark else MUTED)

    def title(self, s, txt, y=0.95, size=32):
        self.text(s, txt, M, y, W - 2 * M, 0.95, size=size, font=HEAD,
                  bold=True, color=WHITE if s._posb_dark else INK)

    def foot(self, s, txt, y=6.45):
        self.text(s, txt, M, y, W - 2 * M, 0.4, size=11.5, italic=True,
                  color=SILVER if s._posb_dark else MUTED)

    def sources(self, s, pairs, y=6.62):
        """Source lines, each attached to the claim it actually supports.

        A single citation in the footer reads as "the source for this slide",
        and stops being true the moment a second source appears on it -- which
        is exactly what happened when a movie from one paper landed under a
        number from another. Taking (what, citation) pairs makes the omission
        visible while you are writing the slide rather than while someone is
        looking at it.

            d.sources(s, [("D = 7.7 µm²/s", "Elowitz et al. 1999"),
                          ("movie", "Valverde-Mendez et al. 2025, Movie S4")])
        """
        dark = s._posb_dark
        wid = (W - 2 * M) / len(pairs)
        for i, (what, cite) in enumerate(pairs):
            x = M + i * wid
            # NOT .upper(): these labels carry units, and case is meaningful
            # in a unit. "µm²/s" upper-cases to "MM²/S", which is a different
            # quantity and a thousand times bigger.
            self.text(s, what, x, y, wid - 0.25, 0.22, size=9,
                      bold=True, color=CYAN if dark else TEAL)
            self.text(s, cite, x, y + 0.21, wid - 0.25, 0.4, size=9.5,
                      italic=True, color=SILVER if dark else MUTED)
        return y + 0.6

    def notes(self, s, txt):
        s.notes_slide.notes_text_frame.text = txt

    # -- the paper-figure slot ----------------------------------------------
    def unattributed(self, s, x, y, w, note="source not yet identified"):
        """A visible ATTRIBUTION NEEDED stamp, on the slide, not in the notes.

        A speaker note is invisible in presentation mode, so a figure whose
        source is unresolved can be projected in front of eighty people with
        nothing to stop it. This is the stop: an amber bar that is impossible to
        miss while rehearsing and embarrassing enough to remove before delivery.

        The intended lifecycle is that this call gets deleted, not that it gets
        tolerated.
        """
        self.shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.42,
                   fill=None, line=AMBER, lw=2)
        self.text(s, "ATTRIBUTION NEEDED", x + 0.15, y + 0.06, 1.75, 0.3,
                  size=10, bold=True, color=AMBER)
        self.text(s, note, x + 1.95, y + 0.07, w - 2.1, 0.3, size=11,
                  italic=True, color=AMBER)
        self.unattributed_figures.append(note)
        return y + 0.42

    def paper_figure(self, s, key, x, y, w, h, ref, caption):
        """Embed a figure from a published paper -- if it is available locally.

        Copyrighted figures cannot live in this public CC BY repository. So the
        image is looked up in `private/paper-figures/<key>.png`, which is
        gitignored:

          * present  -> the real figure is embedded (your classroom deck)
          * absent   -> a labelled slot is drawn naming the exact figure
                        (what CI builds, and what a fork sees)

        Same source, same command, two outputs. Nothing is ever edited by hand.
        """
        img = _resolve_paper_figure(key)
        if img.exists():
            ix, iy, iw, ih = self.image(s, img, x, y, w, h)   # records img
            self.text(s, ref, ix, iy + ih + 0.04, iw, 0.25, size=9,
                      italic=True, color=MUTED if not s._posb_dark else SILVER,
                      align="c")
            # A slot whose aspect ratio does not match the figure's wastes the
            # space it reserved: the image shrinks to fit and lands small in a
            # big empty box. Cheap to detect, invisible until you look at the
            # rendered slide, so detect it.
            fill = (iw * ih) / (w * h)
            if fill < 0.75:
                self.loose_slots.append(
                    (key, round(fill, 2), (round(w, 2), round(h, 2)),
                     (round(iw, 2), round(ih, 2))))
            return True

        self.missing_figures.append((key, ref))
        box = self.shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
                         fill=CARD, line=AMBER, lw=2)
        box.line.dash_style = 4      # dashed
        self.text(s, "FIGURE FROM THE PAPER", x, y + h / 2 - 0.55, w, 0.3,
                  size=9.5, bold=True, color=AMBER, align="c")
        self.text(s, ref, x, y + h / 2 - 0.22, w, 0.35, size=14, font=HEAD,
                  bold=True, color=INK, align="c")
        self.text(s, caption, x + 0.2, y + h / 2 + 0.14, w - 0.4, 0.6,
                  size=11.5, italic=True, color=MUTED, align="c")
        return False

    # -- the reading handed out at the end of this class ---------------------
    def assigned_here(self):
        """Papers this session hands out: the ones the NEXT session discusses.

        Read from readings.yaml through the same resolver the checker uses, so
        the slide and docs/readings.md cannot say different things.
        """
        from tools.build_readings import resolve
        rows, errors = resolve(readings_spec(), sessions())
        if errors:
            raise RuntimeError(
                "readings.yaml does not validate, so no deck can be built from "
                "it. Run `python tools/build_readings.py` to see why.")
        return [r for r in rows if r["assign"] == self.session]

    def discussed_here(self):
        """Papers this session discusses -- i.e. what they were told to read."""
        from tools.build_readings import resolve
        rows, _ = resolve(readings_spec(), sessions())
        return [r for r in rows if r["discuss"] == self.session]

    def assigned_on(self, x, y, w, s, prefix="You read this for today"):
        """A one-line reminder of when the reading for this session went out.

        Small, and it earns its place: it is the visible half of the contract.
        Students who see the date on the slide learn that the assignment is
        real, which is most of why they do it the following week.
        """
        rows = self.discussed_here()
        if not rows:
            return 0
        when = sessions()[rows[0]["assign"]]["date"].strftime("%A %d %B").replace(" 0", " ")
        self.text(s, f"{prefix} — assigned {when}.", x, y, w, 0.3, size=11,
                  italic=True, color=SILVER if s._posb_dark else MUTED)
        return len(rows)

    def assignment(self, s, y=5.55, label="Before next class"):
        """Render the reading due at the next meeting. Returns the box bottom.

        Returning the bottom rather than a count so callers can put the next
        thing under it: the box grows with the number of papers, and a hard-coded
        y below it is a collision waiting for the second reading to be added.

        Nothing is typed here. If readings.yaml says session N+1 discusses a
        paper, the slide says so; if the paper moves, the slide moves with it.
        A deck that quietly disagrees with the syllabus is how a class ends up
        arguing about a paper half the room has not read.
        """
        rows = self.assigned_here()
        self.assignment_rendered = True
        if not rows:
            return y

        # Required first, optional last -- the order they should be done in,
        # not alphabetical order, which is what the file happens to be sorted by.
        rows = sorted(rows, key=lambda r: (not r["required"], r.get("key", "")))

        h = 0.52 + 0.56 * len(rows)
        # The box grows with the number of readings, and a deck written when
        # there was one paper will happily run it off the bottom when a second
        # is added. Clamp, and tell the build so the deck gets reflowed rather
        # than shipped with a box hanging over the edge.
        if y + h > H - 0.18:
            self.assignment_overflow = (self.session, round(y, 2),
                                        round(y + h - (H - 0.18), 2))
            y = H - 0.18 - h
        dark = s._posb_dark
        nxt = sessions().get(self.session + 1)
        when = nxt["date"].strftime("%A").rstrip() if nxt else "next class"
        self.shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, y, W - 2 * M, h,
                   fill=None if dark else WASH, line=CYAN if dark else TEAL, lw=1.5)
        self.text(s, f"{label.upper()} — {when.upper()}", M + 0.3, y + 0.1,
                  6.0, 0.26, size=10, bold=True, color=CYAN if dark else TEAL)
        for i, r in enumerate(rows):
            yy = y + 0.4 + i * 0.56
            tag = "READ" if r["required"] else "OPTIONAL"
            self.text(s, tag, M + 0.3, yy + 0.02, 0.95, 0.24, size=9.5,
                      bold=True, color=AMBER if r["required"] else MUTED)
            self.text(s, short_cite(r), M + 1.3, yy, W - 2 * M - 1.7, 0.28,
                      size=12.5, bold=True, color=WHITE if dark else INK)
            self.text(s, first_sentence(r.get("focus")), M + 1.3, yy + 0.25,
                      W - 2 * M - 1.7, 0.28, size=10.5,
                      italic=True, color=MINT if dark else BODY)
        return y + h

    # -- video ---------------------------------------------------------------
    def movie(self, s, path, x, y, w, h, poster=None, caption=None):
        """Embed a video, letterboxed into the box and centred.

        PowerPoint plays an embedded mp4 in presentation mode; the poster frame
        is what shows until you click it, so it should be the END state of the
        clip rather than the first frame -- a still of the finished picture
        reads as a figure if the video never gets played, which happens more
        often than anyone admits.
        """
        from pptx.util import Inches as _In
        p = self._used(path)
        post = Path(poster) if poster else p.with_name(p.stem + "_poster.png")
        if not post.is_absolute():
            post = ROOT / post
        if not post.exists():
            post = _extract_poster(p, post)
        self._used(post)

        if post.exists():
            from PIL import Image as _Image
            iw, ih = _Image.open(post).size
            scale = min(w / iw, h / ih)
            dw, dh = iw * scale, ih * scale
            x += (w - dw) / 2
            y += (h - dh) / 2
            w, h = dw, dh

        s.shapes.add_movie(str(p), _In(x), _In(y), _In(w), _In(h),
                           poster_frame_image=str(post) if post.exists() else None,
                           mime_type="video/mp4")
        if caption:
            self.text(s, caption, x, y + h + 0.04, w, 0.25, size=9,
                      italic=True, align="c",
                      color=SILVER if s._posb_dark else MUTED)
        return y + h

    def paper_movie(self, s, key, x, y, w, h, ref, caption):
        """A movie from a published paper -- if it is available locally.

        Exactly the paper_figure contract, for video. Present in
        private/paper-movies/<key>.mp4 -> embedded. Absent -> a labelled slot
        naming the movie and where to get it, so the public build and a fork
        see an instruction rather than a hole.

        Supplementary movies are usually the one part of a paper that cannot be
        replaced by a still, which is exactly why they are worth the trouble.
        """
        mp4 = PAPER_MOVIES / f"{key}.mp4"
        if mp4.exists():
            bottom = self.movie(s, mp4, x, y, w, h, caption=ref)
            return True

        self.missing_movies.append((key, ref))
        box = self.shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
                         fill=CARD, line=AMBER, lw=2)
        box.line.dash_style = 4
        self.text(s, "▶  MOVIE FROM THE PAPER", x, y + h / 2 - 0.6, w, 0.3,
                  size=9.5, bold=True, color=AMBER, align="c")
        self.text(s, ref, x, y + h / 2 - 0.26, w, 0.35, size=14, font=HEAD,
                  bold=True, color=INK, align="c")
        self.text(s, caption, x + 0.2, y + h / 2 + 0.12, w - 0.4, 0.7,
                  size=11.5, italic=True, color=MUTED, align="c")
        self.text(s, f"private/paper-movies/{key}.mp4", x, y + h - 0.28, w, 0.25,
                  size=9, italic=True, color=MUTED, align="c")
        return False

    # -- pacing --------------------------------------------------------------
    def pacing(self, total=89, thin=4.5):
        """Minutes of exposition per slide, per segment.

        The failure this catches: a segment where you talk for eight minutes
        with nothing on the screen changing. That is not a short lecture, it is
        an improvised one, and it is invisible in the source -- the deck looks
        fine, the slide looks fine, and only the ratio gives it away.

        Reference point: the 2025 decks ran about 2.2 minutes per slide, with
        almost no student working time. This course gives 35-48% of the period
        back to the students, so the exposition slide count halves -- but the
        RATE should not change much. Anything past `thin` minutes on one
        exposition slide is flagged.

        Segments whose label names an activity are excluded: during a vote or a
        faded worked set the slide is static on purpose. So are segments whose
        label says "board" -- see BOARD_WORDS -- but those are reported
        separately, because "I will derive this at the board" is a claim about
        how the session runs and should not hide inside the student-time number.

        Returns (rows, summary).
        """
        rows, expo_min, expo_slides, act_min, board_min = [], 0, 0, 0, 0
        merged, last = [], None
        for badge, label in self.segments:
            if badge == last and merged:
                merged[-1][2] += 1
            else:
                merged.append([badge, label, 1])
                last = badge

        for badge, label, k in merged:
            nums = re.findall(r"\d+", badge)
            mins = int(nums[1]) - int(nums[0]) if len(nums) >= 2 else 0
            low = label.lower()
            activity = any(w in low for w in self.ACTIVITY_WORDS)
            board = (not activity) and any(w in low for w in self.BOARD_WORDS)
            if activity:
                act_min += mins
            elif board:
                board_min += mins
            else:
                expo_min += mins
                expo_slides += k
            rows.append({"badge": badge, "label": label, "slides": k,
                         "minutes": mins, "activity": activity, "board": board,
                         "per_slide": mins / k if k else 0,
                         "thin": (not activity) and (not board)
                                 and k and mins / k > thin})
        rate = expo_min / expo_slides if expo_slides else 0
        return rows, {"exposition_min": expo_min, "exposition_slides": expo_slides,
                      "min_per_slide": rate, "activity_min": act_min,
                      "activity_frac": act_min / total if total else 0,
                      "board_min": board_min,
                      "thin": [r for r in rows if r["thin"]]}

    # -- output --------------------------------------------------------------
    def save(self, path):
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        return path
